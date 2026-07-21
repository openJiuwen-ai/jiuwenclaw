"""Generate JiuwenAvatar detailed design PPT from design doc structure."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Theme colors
PRIMARY = RGBColor(0x1A, 0x56, 0xDB)  # blue
DARK = RGBColor(0x1E, 0x29, 0x3B)
ACCENT = RGBColor(0x0E, 0xA5, 0xE9)
TEXT = RGBColor(0x33, 0x41, 0x55)
LIGHT = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "zh" / "JiuwenAvatar-详细设计.pptx"

SLIDES = [
    {
        "title": "背景：从对话平台到数字分身任务平台",
        "bullets": [
            "JiuwenSwarm 已支持：多端接入、Gateway 路由、Agent 执行、Cron/Heartbeat、Skill/Memory/Team",
            "原有模式：用户发起请求 → Gateway 转发 → AgentServer 执行 → 返回结果",
            "数字分身场景缺少四项能力：",
            "  · 不知道哪个分身应被自动唤醒",
            "  · 自动执行无法记录与追踪",
            "  · 长任务无法取消、失败无法沉淀",
            "  · 执行结果缺少结构化报告与已读/统计",
            "JiuwenAvatar 目标：从「Agent 对话平台」升级为「数字分身任务执行平台」",
        ],
    },
    {
        "title": "0 层业务闭环",
        "bullets": [
            "核心闭环：Persona → Avatar → Trigger → Mission → Agent Execution → Report",
            "Persona：角色模板（Committer / Developer / Tester 等）",
            "Avatar：用户创建的数字分身实例",
            "Trigger：定义何时唤醒哪个 Avatar",
            "Mission：记录一次自动任务执行",
            "Agent Execution：AgentServer 按 Avatar 身份执行任务",
            "Report：任务结果的结构化沉淀",
            "新增设计构建端到端自动任务闭环，而非单点加功能",
        ],
    },
    {
        "title": "系统分工：控制面与执行面",
        "bullets": [
            "Gateway 进程（控制面）：",
            "  · 触发、派发、Mission 生命周期、Report、前端 API、结果推送",
            "AgentServer 进程（执行面）：",
            "  · 解析 avatar_id、注入 Persona/Skill/CodingEngine 上下文、调用 Agent 执行",
            "关键边界：",
            "  · Gateway 不理解 Agent 内部执行细节，只负责控制任务生命周期",
            "  · AgentServer 不持久化 Mission/Report，只负责以正确 Avatar 上下文执行",
            "  · avatar_id 是控制面与执行面的核心关联字段",
            "  · session_id 是 Mission 取消时关联 AgentServer 会话的关键字段",
        ],
    },
    {
        "title": "领域模型",
        "bullets": [
            "Persona（模板）：id, skills, coding_engines, trigger_templates, report_template, system_prompt",
            "Avatar（实例）：persona_id, skills, coding_engine, trigger_ids, report_channels, status",
            "TriggerConfig：type(cron/heartbeat/webhook/event), avatar_id, trigger_prompt, generate_report",
            "Mission（账本）：avatar_id, trigger_id, status(pending/running/completed/failed/cancelled), session_id",
            "MissionReport（输出）：mission_id, title, summary, sections, metrics",
            "CodingEngine：jiuwen-coding / claude-code / codex，统一 ensure_ready() + run_task() 契约",
            "关系：Persona → Avatar → Trigger → Mission → MissionReport",
        ],
    },
    {
        "title": "创建链路：Avatar 与 Trigger 自动 Provision",
        "bullets": [
            "1. 前端选择 Persona 创建 Avatar（avatars.create）",
            "2. PersonaManager 合并 Persona 默认技能与用户额外技能",
            "3. 解析并保存 coding_engine",
            "4. AvatarFactory 安装 Avatar 绑定的内置 Skill",
            "5. 根据 Persona 的 trigger_templates 自动创建 Trigger",
            "6. 保存 trigger_ids 到 Avatar，返回前端",
            "关键代码：persona/manager.py, avatar_factory.py, gateway/trigger/engine.py",
        ],
    },
    {
        "title": "触发链路：ITrigger.fire → _dispatch_fire",
        "bullets": [
            "ITrigger.fire() 是所有触发器进入执行闭环的统一入口",
            "Cron / Heartbeat / Webhook / Event 只负责「什么时候触发」",
            "TriggerEngine._dispatch_fire 统一负责「触发后如何执行」：",
            "  1. 更新 last_triggered_at，生成 run_id / session_id",
            "  2. MissionManager.create_mission() → 状态 RUNNING",
            "  3. AgentClient.send_request(E2A chat.send, params.avatar_id + prompt)",
            "  4. 提取 result_text → Mission COMPLETED / FAILED",
            "  5. 若 generate_report → create_report()",
            "  6. 若 target_channel → publish_robot_messages() 回推结果",
        ],
    },
    {
        "title": "执行链路：avatar_id 注入分身上下文",
        "bullets": [
            "Gateway 发送 E2A chat.send(params.avatar_id, prompt)",
            "AgentAdapter 提取 avatar_id → resolve_avatar_chat_context()",
            "PersonaAvatarChatRail.set_context() 注入身份约束",
            "ensure_avatar_skills_installed() 安装缺失 Skill",
            "get_coding_engine() → ensure_ready() → 注册或移除 coding_task",
            "SkillRail 白名单限制可用 Skill",
            "Runner / DeepAgent 执行任务并返回结果",
            "同一 AgentServer 可处理普通对话与指定 Avatar 的自动任务，由 avatar_id 决定",
        ],
    },
    {
        "title": "报告链路：Mission 生命周期与 MissionReport",
        "bullets": [
            "Mission 状态：pending → running → completed / failed / cancelled",
            "持久化：~/.jiuwenavatar/reports/missions.json",
            "MissionReport 在任务成功后生成：title, summary, sections, metrics",
            "持久化：~/.jiuwenavatar/reports/reports.json",
            "ReadState：记录任务和报告已读状态（read_state.json）",
            "UsageStats：累计使用统计（active_days, total_tasks, completed_tasks 等）",
            "Web API：missions.list/get/stats, reports.list/get, report.unread_counts",
        ],
    },
    {
        "title": "取消链路：Mission Cancel 与 CHAT_CANCEL",
        "bullets": [
            "取消入口：前端 missions.cancel → TriggerEngine.cancel_mission()",
            "Gateway 读取 Mission.session_id",
            "向 AgentServer 发送 E2A CHAT_CANCEL 取消会话",
            "MissionManager.cancel_mission() 将状态置为 CANCELLED",
            "状态保护：Mission 已 CANCELLED 后，后续 COMPLETED/FAILED 不会覆盖",
            "避免用户取消后，AgentServer 稍后返回结果又把 Mission 改成 completed",
        ],
    },
    {
        "title": "AgentServer 新增设计",
        "bullets": [
            "PersonaAvatarChatRail：模型调用前注入 Avatar 身份、Skill 白名单、编码约束",
            "  · 禁止 search_skill / install_skill / browser_agent 等绕行方式",
            "  · 外部 CLI 编码后端必须通过 coding_task 委派",
            "AvatarChatContext：解析 avatar_id → system_prompt, skills, coding_engine",
            "CodingEngine 统一抽象：",
            "  · jiuwen-coding：Leader 直接用 Skill + bash，不注册 coding_task",
            "  · claude-code / codex：注册 coding_task，CLI 工作区按 avatar_id 隔离",
            "AgentAdapter：动态注册/移除 coding_task，设置 Skill 白名单",
        ],
    },
    {
        "title": "与 JiuwenSwarm 的差异",
        "bullets": [
            "JiuwenSwarm 原有：Gateway, MessageHandler, Cron/Heartbeat, AgentServer, Session/Skill/A2UI",
            "JiuwenAvatar 新增/增强：",
            "  · TriggerEngine + WebhookTrigger + EventTrigger",
            "  · Trigger 与 Avatar 绑定 + trigger_templates 自动 provisioning",
            "  · Mission 生命周期 + MissionReport + ReadState + UsageStats",
            "  · Mission cancel 与 AgentServer CHAT_CANCEL 打通",
            "  · Persona/Avatar 运行时上下文 + PersonaAvatarChatRail",
            "  · CodingEngine 统一抽象 + coding_task + 按 Avatar 隔离 CLI 工作区",
            "核心变化：可追踪 Mission → 指定 Avatar 身份执行 → 生成 Report → 前端查询/取消/统计",
        ],
    },
    {
        "title": "风险与演进方向",
        "bullets": [
            "当前风险：",
            "  · Trigger/Mission/Report 使用 JSON 文件存储，并发写入能力有限",
            "  · Webhook 可无 secret，生产环境存在安全风险",
            "  · Gateway 单进程调度，分布式场景需调度锁或 leader election",
            "  · CLI 编码引擎依赖本机环境、凭据和安装状态",
            "后续优化：",
            "  · 存储升级为 SQLite 或服务端数据库",
            "  · Webhook 增加 GitCode、飞书等平台级 payload 解析",
            "  · Report 引入 Persona report_template",
            "  · Mission 重试/超时/失败归因 + Trigger 执行审计日志",
            "  · CodingEngine 前端健康检查 + 自动化测试覆盖",
        ],
    },
]


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_bar(slide, prs: Presentation, title: str) -> None:
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0),
        Inches(0),
        prs.slide_width,
        Inches(1.1),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    accent = slide.shapes.add_shape(
        1,
        Inches(0),
        Inches(1.1),
        prs.slide_width,
        Inches(0.06),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(12), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Microsoft YaHei"


def _add_bullets(slide, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(12.0), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(18 if not text.startswith("  ") else 16)
        p.font.color.rgb = TEXT if not text.startswith("  ") else LIGHT
        p.space_after = Pt(6)
        p.level = 1 if text.startswith("  ") else 0


def _add_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(2.8), prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.2))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = "JiuwenAvatar 自动任务闭环"
    tp.font.size = Pt(40)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.font.name = "Microsoft YaHei"
    tp.alignment = PP_ALIGN.LEFT

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.8))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = "详细设计"
    sp.font.size = Pt(28)
    sp.font.color.rgb = ACCENT
    sp.font.name = "Microsoft YaHei"

    info_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(2.5))
    tf = info_box.text_frame
    for i, line in enumerate(
        [
            "从 Agent 对话平台到数字分身任务执行平台",
            "Gateway 控制面 + AgentServer 执行面",
            "Persona → Avatar → Trigger → Mission → Report",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = LIGHT
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(10)


def _add_content_slide(prs: Presentation, data: dict, index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, prs, data["title"])
    _add_bullets(slide, data["bullets"])

    page = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4))
    pp = page.text_frame.paragraphs[0]
    pp.text = str(index)
    pp.font.size = Pt(12)
    pp.font.color.rgb = LIGHT
    pp.font.name = "Microsoft YaHei"
    pp.alignment = PP_ALIGN.RIGHT


def _add_closing(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, PRIMARY)

    box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(2))
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.text = "谢谢"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.font.name = "Microsoft YaHei"
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "JiuwenAvatar · 数字分身任务执行平台"
    p2.font.size = Pt(20)
    p2.font.color.rgb = WHITE
    p2.font.name = "Microsoft YaHei"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_cover(prs)
    for i, slide_data in enumerate(SLIDES, start=1):
        _add_content_slide(prs, slide_data, i)
    _add_closing(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build_presentation()
    print(f"PPT generated: {path}")
