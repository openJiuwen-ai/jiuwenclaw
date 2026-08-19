import os
from pathlib import Path

from openjiuwen.core.foundation.tool import Tool, ToolCard

from pdf_font_utils import select_pdf_font


class FormatConverter(Tool):
    """格式互转工具：支持 PDF、Word、Excel、CSV、PPT、Markdown 之间的格式转换。"""

    def __init__(self) -> None:
        super().__init__(
            ToolCard(
                id="format_converter",
                name="format_converter",
                description=(
                    "格式互转工具：支持PDF↔Word、Word→PDF、Excel↔CSV、"
                    "PPT→PDF、Markdown→Word/PPT等格式转换。"
                    "当用户需要转换文档格式时调用。"
                ),
                input_params={
                    "type": "object",
                    "properties": {
                        "source_path": {
                            "type": "string",
                            "description": "源文件绝对路径",
                        },
                        "target_format": {
                            "type": "string",
                            "enum": ["word", "pdf", "excel", "csv", "ppt"],
                            "description": "目标格式",
                        },
                        "output_subdir": {
                            "type": "string",
                            "description": "输出子目录名，默认为 converted",
                        },
                    },
                    "required": ["source_path", "target_format"],
                },
            )
        )

    async def invoke(self, inputs, **kwargs):
        source_path = inputs.get("source_path", "")
        target_format = inputs.get("target_format", "")
        output_subdir = inputs.get("output_subdir", "converted")

        if not source_path or not os.path.isfile(source_path):
            return {"success": False, "error": f"源文件不存在: {source_path}"}
        if not target_format:
            return {"success": False, "error": "缺少 target_format 参数"}

        from openjiuwen.core.sys_operation.cwd import get_cwd

        base_dir = Path(get_cwd()) / output_subdir
        base_dir.mkdir(parents=True, exist_ok=True)

        source_ext = Path(source_path).suffix.lower()
        stem = Path(source_path).stem
        ext_map = {
            "word": ".docx",
            "pdf": ".pdf",
            "excel": ".xlsx",
            "csv": ".csv",
            "ppt": ".pptx",
        }
        output_path = base_dir / (
            f"{stem}_converted{ext_map.get(target_format, f'.{target_format}')}"
        )

        try:
            source_type = self._detect_type(source_ext)
            if source_type is None:
                return {
                    "success": False,
                    "error": f"无法识别源文件类型: {source_ext}",
                }

            if source_type == target_format:
                return {
                    "success": False,
                    "error": "源格式与目标格式相同，无需转换",
                }

            converter = self._get_converter(source_type, target_format)
            if converter is None:
                return {
                    "success": False,
                    "error": f"不支持的转换路径: {source_type} -> {target_format}",
                }

            converter(str(output_path), source_path)

            if not output_path.exists() or output_path.stat().st_size == 0:
                return {
                    "success": False,
                    "error": "转换失败，输出文件为空或不存在",
                }

            return {
                "success": True,
                "source_format": source_type,
                "target_format": target_format,
                "source_path": source_path,
                "path": str(output_path),
                "absolute_path": str(output_path.resolve()),
                "filename": output_path.name,
                "size_bytes": output_path.stat().st_size,
                "exists": True,
            }
        except ImportError as e:
            return {
                "success": False,
                "error": f"依赖库缺失: {e}. 请安装对应依赖后重试。",
            }
        except Exception as e:
            return {"success": False, "error": f"转换失败: {e}"}

    async def stream(self, inputs, **kwargs):
        yield await self.invoke(inputs, **kwargs)

    @staticmethod
    def _detect_type(ext: str) -> str | None:
        if ext == ".pdf":
            return "pdf"
        if ext in (".doc", ".docx"):
            return "word"
        if ext in (".xls", ".xlsx"):
            return "excel"
        if ext == ".csv":
            return "csv"
        if ext in (".ppt", ".pptx"):
            return "ppt"
        if ext in (".md", ".markdown"):
            return "markdown"
        return None

    def _get_converter(self, source_type: str, target_format: str):
        """Return the conversion function or None."""
        mapping = {
            ("pdf", "word"): self._pdf_to_word,
            ("pdf", "excel"): self._pdf_to_excel,
            ("word", "pdf"): self._word_to_pdf,
            ("excel", "csv"): self._excel_to_csv,
            ("csv", "excel"): self._csv_to_excel,
            ("ppt", "pdf"): self._ppt_to_pdf,
            ("markdown", "word"): self._markdown_to_word,
            ("markdown", "ppt"): self._markdown_to_ppt,
        }
        return mapping.get((source_type, target_format))

    @staticmethod
    def _pdf_to_word(output_path: str, source_path: str) -> None:
        from docx import Document
        from pypdf import PdfReader

        reader = PdfReader(source_path)
        doc = Document()
        for page in reader.pages:
            text = page.extract_text() or ""
            if text.strip():
                doc.add_paragraph(text)
                doc.add_page_break()
        doc.save(output_path)

    @staticmethod
    def _word_to_pdf(output_path: str, source_path: str) -> None:
        from docx import Document
        from fpdf import FPDF

        doc = Document(source_path)
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)

        font_name = select_pdf_font(pdf)

        pdf.set_font(font_name, "", 11)
        for para in doc.paragraphs:
            if para.text.strip():
                pdf.multi_cell(0, 7, para.text)
                pdf.ln(3)
        pdf.output(output_path)

    @staticmethod
    def _excel_to_csv(output_path: str, source_path: str) -> None:
        import csv

        from openpyxl import load_workbook

        wb = load_workbook(source_path, data_only=True)
        ws = wb.active
        with open(output_path, "w", encoding="utf-8", newline="") as f:
            writer = csv.writer(f)
            for row in ws.iter_rows(values_only=True):
                writer.writerow(
                    [str(c) if c is not None else "" for c in row]
                )

    @staticmethod
    def _csv_to_excel(output_path: str, source_path: str) -> None:
        import csv

        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        with open(source_path, encoding="utf-8") as f:
            reader = csv.reader(f)
            for row in reader:
                ws.append(row)
        wb.save(output_path)

    @staticmethod
    def _ppt_to_pdf(output_path: str, source_path: str) -> None:
        from fpdf import FPDF
        from pptx import Presentation

        prs = Presentation(source_path)
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)

        font_name = select_pdf_font(pdf)

        for slide in prs.slides:
            pdf.add_page()
            pdf.set_font(font_name, "B", 14)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if text.strip():
                        pdf.multi_cell(0, 7, text)
                        pdf.ln(3)
        pdf.output(output_path)

    @staticmethod
    def _markdown_to_word(output_path: str, source_path: str) -> None:
        import re

        from docx import Document

        with open(source_path, encoding="utf-8") as f:
            lines = f.readlines()

        doc = Document()
        for line in lines:
            line = line.rstrip("\n")
            if not line.strip():
                continue
            m = re.match(r"^(#{1,6})\s+(.+)$", line)
            if m:
                level = len(m.group(1))
                doc.add_heading(m.group(2), level=level)
                continue
            m = re.match(r"^[-*]\s+(.+)$", line)
            if m:
                doc.add_paragraph(m.group(1), style="List Bullet")
                continue
            m = re.match(r"^\d+\.\s+(.+)$", line)
            if m:
                doc.add_paragraph(m.group(1), style="List Number")
                continue
            if line.startswith("|"):
                doc.add_paragraph(line)
                continue
            doc.add_paragraph(line)
        doc.save(output_path)

    @staticmethod
    def _markdown_to_ppt(output_path: str, source_path: str) -> None:
        import re

        from pptx import Presentation

        with open(source_path, encoding="utf-8") as f:
            lines = f.readlines()

        prs = Presentation()
        current_title = ""
        current_body: list[str] = []

        def add_slide_if_ready() -> None:
            nonlocal current_title, current_body
            if not current_title and not current_body:
                return
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            if current_title and slide.shapes.title:
                slide.shapes.title.text = current_title
            if current_body and len(slide.placeholders) > 1:
                slide.placeholders[1].text_frame.text = "\n".join(
                    current_body
                )
            current_title = ""
            current_body = []

        for line in lines:
            line = line.rstrip("\n")
            m = re.match(r"^(#{1,2})\s+(.+)$", line)
            if m:
                add_slide_if_ready()
                current_title = m.group(2)
                continue
            if line.strip():
                current_body.append(line)
        add_slide_if_ready()
        prs.save(output_path)

    @staticmethod
    def _pdf_to_excel(output_path: str, source_path: str) -> None:
        from openpyxl import Workbook
        from pdfplumber import open as pdf_open

        wb = Workbook()
        ws = wb.active
        ws.title = "Extracted"

        with pdf_open(source_path) as pdf:
            for page in pdf.pages:
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        for row in table:
                            ws.append(
                                [
                                    str(c) if c else ""
                                    for c in row
                                ]
                            )
                        ws.append([])
                else:
                    text = page.extract_text() or ""
                    for line in text.split("\n"):
                        ws.append([line])

        wb.save(output_path)
