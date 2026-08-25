import os
from pathlib import Path

from openjiuwen.core.foundation.tool import Tool, ToolCard


def _append_to_run(run, element) -> None:
    getattr(run, "_element").append(element)


class ReportStyler(Tool):
    """高级排版/报告美化工具：Word自动目录、页眉页脚页码、图表插入、统一主题样式。"""

    def __init__(self) -> None:
        super().__init__(
            ToolCard(
                id="report_styler",
                name="report_styler",
                description=(
                    "高级排版/报告美化工具：Word自动目录、页眉页脚页码、"
                    "把数据画成图表插进Word/PPT、统一主题样式。"
                    "当用户需要美化文档排版时调用。"
                ),
                input_params={
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "要美化的文件绝对路径",
                        },
                        "format": {
                            "type": "string",
                            "enum": ["word", "ppt"],
                            "description": "文档格式",
                        },
                        "operations": {
                            "type": "array",
                            "items": {"type": "object"},
                            "description": (
                                "美化操作列表。Word: add_toc/"
                                "add_header_footer/insert_chart/"
                                "apply_theme/set_margins/add_cover_page。"
                                "PPT: apply_theme/insert_chart/"
                                "set_slide_size"
                            ),
                        },
                    },
                    "required": ["file_path", "format", "operations"],
                },
            )
        )

    async def invoke(self, inputs, **kwargs):
        file_path = inputs.get("file_path", "")
        fmt = inputs.get("format", "")
        operations = inputs.get("operations", [])

        if not file_path or not os.path.isfile(file_path):
            return {"success": False, "error": f"文件不存在: {file_path}"}
        if not fmt:
            return {"success": False, "error": "缺少 format 参数"}
        if not operations:
            return {"success": False, "error": "缺少 operations 参数"}

        try:
            # 图表是就地美化的中间产物，跟随被美化的文档存放。
            charts_dir = Path(file_path).resolve().parent / "charts"
            charts_dir.mkdir(parents=True, exist_ok=True)
            if fmt == "word":
                results = self._style_word(file_path, operations, charts_dir)
            elif fmt == "ppt":
                results = self._style_ppt(file_path, operations, charts_dir)
            else:
                return {"success": False, "error": f"不支持的格式: {fmt}"}

            return {
                "success": True,
                "format": fmt,
                "file_path": file_path,
                "operations_executed": len(results),
                "results": results,
                "size_bytes": os.path.getsize(file_path),
                "exists": True,
            }
        except ImportError as e:
            return {
                "success": False,
                "error": f"依赖库缺失: {e}. 请安装对应依赖后重试。",
            }
        except Exception as e:
            return {"success": False, "error": f"美化失败: {e}"}

    async def stream(self, inputs, **kwargs):
        yield await self.invoke(inputs, **kwargs)

    @staticmethod
    def _style_word(
        file_path: str, operations: list, charts_dir: Path
    ) -> list:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        from docx.shared import Inches, Pt, RGBColor

        doc = Document(file_path)
        results = []

        for op in operations:
            op_type = op.get("type", "")

            if op_type == "add_toc":
                toc_para = (
                    doc.paragraphs[0]
                    if doc.paragraphs
                    else doc.add_paragraph()
                )
                run = toc_para.add_run()
                fld_begin = OxmlElement("w:fldChar")
                fld_begin.set(qn("w:fldCharType"), "begin")
                _append_to_run(run, fld_begin)

                run2 = toc_para.add_run()
                instr = OxmlElement("w:instrText")
                instr.set(qn("xml:space"), "preserve")
                instr.text = ' TOC \\o "1-3" \\h \\z \\u '
                _append_to_run(run2, instr)

                run3 = toc_para.add_run()
                fld_end = OxmlElement("w:fldChar")
                fld_end.set(qn("w:fldCharType"), "end")
                _append_to_run(run3, fld_end)

                results.append(
                    {
                        "type": "add_toc",
                        "status": "TOC域已插入，打开Word后按F9更新",
                    }
                )

            elif op_type == "add_header_footer":
                header_text = op.get("header_text", "")
                footer_text = op.get("footer_text", "")
                add_page_number = op.get("add_page_number", False)

                if header_text:
                    for section in doc.sections:
                        header = section.header
                        header.is_linked_to_previous = False
                        if header.paragraphs:
                            header.paragraphs[0].text = header_text
                        else:
                            header.add_paragraph(header_text)

                if footer_text or add_page_number:
                    for section in doc.sections:
                        footer = section.footer
                        footer.is_linked_to_previous = False
                        if footer_text:
                            if footer.paragraphs:
                                footer.paragraphs[0].text = footer_text
                            else:
                                footer.add_paragraph(footer_text)

                        if add_page_number:
                            p = (
                                footer.paragraphs[0]
                                if footer.paragraphs
                                else footer.add_paragraph()
                            )
                            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            run = p.add_run()
                            fld_begin = OxmlElement("w:fldChar")
                            fld_begin.set(qn("w:fldCharType"), "begin")
                            _append_to_run(run, fld_begin)
                            run2 = p.add_run()
                            instr = OxmlElement("w:instrText")
                            instr.text = " PAGE "
                            _append_to_run(run2, instr)
                            run3 = p.add_run()
                            fld_end = OxmlElement("w:fldChar")
                            fld_end.set(qn("w:fldCharType"), "end")
                            _append_to_run(run3, fld_end)

                results.append(
                    {
                        "type": "add_header_footer",
                        "header": bool(header_text),
                        "footer": bool(footer_text),
                        "page_number": add_page_number,
                    }
                )

            elif op_type == "insert_chart":
                chart_type = op.get("chart_type", "bar")
                chart_data = op.get("chart_data", {})
                title = op.get("title", "")

                chart_path = ReportStyler._generate_chart_image(
                    chart_type, chart_data, title, charts_dir
                )
                if chart_path:
                    doc.add_picture(chart_path, width=Inches(6))
                    results.append(
                        {
                            "type": "insert_chart",
                            "chart_type": chart_type,
                            "image": chart_path,
                        }
                    )
                else:
                    results.append(
                        {
                            "type": "insert_chart",
                            "error": "图表生成失败（需matplotlib）",
                        }
                    )

            elif op_type == "apply_theme":
                font_name = op.get("font_name", "Microsoft YaHei")
                font_size = op.get("font_size", 11)
                heading_color = op.get("heading_color", "1F4E79")

                for para in doc.paragraphs:
                    for run in para.runs:
                        run.font.name = font_name
                        run.font.size = Pt(font_size)

                for para in doc.paragraphs:
                    if para.style and para.style.name.startswith(
                        "Heading"
                    ):
                        for run in para.runs:
                            r = int(heading_color[0:2], 16)
                            g = int(heading_color[2:4], 16)
                            b = int(heading_color[4:6], 16)
                            run.font.color.rgb = RGBColor(r, g, b)

                results.append(
                    {
                        "type": "apply_theme",
                        "font": font_name,
                        "size": font_size,
                        "heading_color": heading_color,
                    }
                )

            elif op_type == "set_margins":
                from docx.shared import Cm

                top = op.get("top", 2.54)
                bottom = op.get("bottom", 2.54)
                left = op.get("left", 3.18)
                right = op.get("right", 3.18)

                for section in doc.sections:
                    section.top_margin = Cm(top)
                    section.bottom_margin = Cm(bottom)
                    section.left_margin = Cm(left)
                    section.right_margin = Cm(right)

                results.append(
                    {
                        "type": "set_margins",
                        "top": top,
                        "bottom": bottom,
                        "left": left,
                        "right": right,
                    }
                )

            elif op_type == "add_cover_page":
                title = op.get("title", "")
                subtitle = op.get("subtitle", "")
                author = op.get("author", "")
                date_str = op.get("date", "")

                body = doc.element.body
                first_child = (
                    body[0] if len(body) > 0 else None
                )

                def insert_para_before(
                    doc_body, ref_elem, text, size_pt=12, bold=False
                ):
                    p = OxmlElement("w:p")
                    para_props = OxmlElement("w:pPr")
                    jc = OxmlElement("w:jc")
                    jc.set(qn("w:val"), "center")
                    para_props.append(jc)
                    p.append(para_props)
                    if text:
                        r = OxmlElement("w:r")
                        run_props = OxmlElement("w:rPr")
                        sz = OxmlElement("w:sz")
                        sz.set(
                            qn("w:val"),
                            str(size_pt * 2),
                        )
                        run_props.append(sz)
                        if bold:
                            b_el = OxmlElement("w:b")
                            run_props.append(b_el)
                        r.append(run_props)
                        t = OxmlElement("w:t")
                        t.text = text
                        r.append(t)
                        p.append(r)
                    if ref_elem is not None:
                        ref_elem.addprevious(p)
                    else:
                        doc_body.append(p)
                    return p

                ref = first_child
                for _ in range(6):
                    ref = insert_para_before(body, ref, "")

                if title:
                    ref = insert_para_before(
                        body, ref, title, size_pt=28, bold=True
                    )
                if subtitle:
                    ref = insert_para_before(
                        body, ref, subtitle, size_pt=16
                    )

                for _ in range(4):
                    ref = insert_para_before(body, ref, "")

                if author:
                    ref = insert_para_before(
                        body, ref, author, size_pt=12
                    )
                if date_str:
                    ref = insert_para_before(
                        body, ref, date_str, size_pt=12
                    )

                p_break = OxmlElement("w:p")
                r_break = OxmlElement("w:r")
                br = OxmlElement("w:br")
                br.set(qn("w:type"), "page")
                r_break.append(br)
                p_break.append(r_break)
                if ref is not None:
                    ref.addnext(p_break)

                results.append(
                    {
                        "type": "add_cover_page",
                        "title": title[:50],
                    }
                )

        doc.save(file_path)
        return results

    @staticmethod
    def _style_ppt(
        file_path: str, operations: list, charts_dir: Path
    ) -> list:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches

        prs = Presentation(file_path)
        results = []

        for op in operations:
            op_type = op.get("type", "")

            if op_type == "apply_theme":
                font_name = op.get("font_name", "Microsoft YaHei")
                title_color = op.get("title_color", "1F4E79")
                body_color = op.get("body_color", "333333")

                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in (
                                shape.text_frame.paragraphs
                            ):
                                for run in para.runs:
                                    run.font.name = font_name
                                    if shape == slide.shapes.title:
                                        r = int(title_color[0:2], 16)
                                        g = int(title_color[2:4], 16)
                                        b = int(title_color[4:6], 16)
                                        run.font.color.rgb = RGBColor(
                                            r, g, b
                                        )
                                    else:
                                        r = int(body_color[0:2], 16)
                                        g = int(body_color[2:4], 16)
                                        b = int(body_color[4:6], 16)
                                        run.font.color.rgb = RGBColor(
                                            r, g, b
                                        )

                results.append(
                    {
                        "type": "apply_theme",
                        "font": font_name,
                        "title_color": title_color,
                        "body_color": body_color,
                    }
                )

            elif op_type == "insert_chart":
                chart_type = op.get("chart_type", "bar")
                chart_data = op.get("chart_data", {})
                title = op.get("title", "")
                slide_index = op.get("slide_index", -1)

                chart_path = ReportStyler._generate_chart_image(
                    chart_type, chart_data, title, charts_dir
                )
                if chart_path:
                    if 0 <= slide_index < len(prs.slides):
                        slide = prs.slides[slide_index]
                    else:
                        slide_layout = prs.slide_layouts[6]
                        slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.add_picture(
                        chart_path,
                        Inches(1),
                        Inches(2),
                        Inches(8),
                        Inches(4),
                    )
                    results.append(
                        {
                            "type": "insert_chart",
                            "chart_type": chart_type,
                            "slide": (
                                slide_index
                                if slide_index >= 0
                                else "new"
                            ),
                        }
                    )
                else:
                    results.append(
                        {
                            "type": "insert_chart",
                            "error": "图表生成失败（需matplotlib）",
                        }
                    )

            elif op_type == "set_slide_size":
                width = op.get("width", 10)
                height = op.get("height", 7.5)
                prs.slide_width = Inches(width)
                prs.slide_height = Inches(height)
                results.append(
                    {
                        "type": "set_slide_size",
                        "width": width,
                        "height": height,
                    }
                )

        prs.save(file_path)
        return results

    @staticmethod
    def _generate_chart_image(
        chart_type: str,
        chart_data: dict,
        title: str,
        charts_dir: Path,
    ) -> str | None:
        """Generate a chart image, return image path or None."""
        try:
            import matplotlib

            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
        except ImportError:
            return None

        labels = chart_data.get("labels", [])
        values = chart_data.get("values", [])
        if not labels or not values:
            return None

        fig, ax = plt.subplots(figsize=(8, 5))

        if chart_type == "bar":
            ax.bar(labels, values)
        elif chart_type == "line":
            ax.plot(labels, values, marker="o")
        elif chart_type == "pie":
            ax.pie(values, labels=labels, autopct="%1.1f%%")
        elif chart_type == "scatter":
            ax.scatter(labels, values)
        else:
            ax.bar(labels, values)

        if title:
            ax.set_title(title)

        plt.rcParams["font.sans-serif"] = [
            "SimHei",
            "Microsoft YaHei",
            "DejaVu Sans",
        ]
        plt.rcParams["axes.unicode_minus"] = False

        img_path = str(
            charts_dir / f"chart_{chart_type}_{abs(hash(title)) % 10000}.png"
        )
        fig.savefig(img_path, dpi=150, bbox_inches="tight")
        plt.close(fig)
        return img_path
