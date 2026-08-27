import os
from pathlib import Path

from openjiuwen.core.foundation.tool import Tool, ToolCard


class DocumentParser(Tool):
    """文档解析工具：支持 PDF、Word、Excel、PPT 文件的内容提取。"""

    def __init__(self) -> None:
        super().__init__(
            ToolCard(
                id="document_parser",
                name="document_parser",
                description=(
                    "文档解析工具：解析PDF、Word、Excel、PPT文件，"
                    "提取文本、表格与结构化内容。"
                    "当用户上传文档需要提取内容、整理信息时调用。"
                ),
                input_params={
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "要解析的文件绝对路径",
                        },
                        "file_type": {
                            "type": "string",
                            "enum": ["pdf", "word", "excel", "ppt", "auto"],
                            "description": "文件类型，auto为根据扩展名自动检测",
                        },
                        "extract_tables": {
                            "type": "boolean",
                            "description": "是否提取表格数据，默认true",
                        },
                    },
                    "required": ["file_path"],
                },
            )
        )

    async def invoke(self, inputs, **kwargs):
        file_path = inputs.get("file_path", "")
        file_type = inputs.get("file_type", "auto")
        extract_tables = inputs.get("extract_tables", True)

        if not file_path or not os.path.isfile(file_path):
            return {"success": False, "error": f"文件不存在: {file_path}"}

        if file_type == "auto":
            file_type = self._detect_type(file_path)
            if file_type is None:
                return {
                    "success": False,
                    "error": "无法识别文件类型，请指定 file_type 参数",
                }

        try:
            if file_type == "pdf":
                result = self._parse_pdf(file_path, extract_tables)
            elif file_type == "word":
                result = self._parse_word(file_path, extract_tables)
            elif file_type == "excel":
                result = self._parse_excel(file_path)
            elif file_type == "ppt":
                result = self._parse_ppt(file_path)
            else:
                return {"success": False, "error": f"不支持的文件类型: {file_type}"}

            return {
                "success": True,
                "file_type": file_type,
                "file_path": file_path,
                "content": result,
            }
        except ImportError as e:
            return {
                "success": False,
                "error": f"依赖库缺失: {e}. 请安装对应依赖后重试。",
            }
        except Exception as e:
            return {"success": False, "error": f"解析失败: {e}"}

    async def stream(self, inputs, **kwargs):
        yield await self.invoke(inputs, **kwargs)

    @staticmethod
    def _detect_type(file_path: str) -> str | None:
        ext = Path(file_path).suffix.lower()
        if ext == ".pdf":
            return "pdf"
        if ext in (".doc", ".docx"):
            return "word"
        if ext in (".xls", ".xlsx", ".csv"):
            return "excel"
        if ext in (".ppt", ".pptx"):
            return "ppt"
        return None

    @staticmethod
    def _parse_pdf(file_path: str, extract_tables: bool) -> dict:
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages.append({"page": i + 1, "text": text})

        tables: list = []
        if extract_tables:
            try:
                import pdfplumber

                with pdfplumber.open(file_path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        page_tables = page.extract_tables()
                        if page_tables:
                            for t_idx, table in enumerate(page_tables):
                                tables.append(
                                    {
                                        "page": i + 1,
                                        "table_index": t_idx,
                                        "data": table,
                                    }
                                )
            except ImportError:
                pass

        return {"total_pages": len(reader.pages), "pages": pages, "tables": tables}

    @staticmethod
    def _parse_word(file_path: str, extract_tables: bool) -> dict:
        from docx import Document

        doc = Document(file_path)
        paragraphs = []
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                paragraphs.append(
                    {
                        "index": i,
                        "text": para.text,
                        "style": para.style.name if para.style else "Normal",
                    }
                )

        tables: list = []
        if extract_tables:
            for t_idx, table in enumerate(doc.tables):
                rows = []
                for row in table.rows:
                    rows.append([cell.text for cell in row.cells])
                tables.append({"table_index": t_idx, "data": rows})

        return {
            "total_paragraphs": len(paragraphs),
            "paragraphs": paragraphs,
            "tables": tables,
        }

    @staticmethod
    def _parse_excel(file_path: str) -> dict:
        ext = Path(file_path).suffix.lower()
        if ext == ".csv":
            import csv

            rows = []
            with open(file_path, encoding="utf-8") as f:
                reader = csv.reader(f)
                for row in reader:
                    rows.append([str(cell) for cell in row])
            max_col = max(len(r) for r in rows) if rows else 0
            return {
                "sheets": [
                    {
                        "sheet_name": "Sheet1",
                        "rows": rows,
                        "max_row": len(rows),
                        "max_col": max_col,
                    }
                ]
            }

        from openpyxl import load_workbook

        wb = load_workbook(file_path, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append(
                    [str(cell) if cell is not None else "" for cell in row]
                )
            sheets.append(
                {
                    "sheet_name": sheet_name,
                    "rows": rows,
                    "max_row": ws.max_row,
                    "max_col": ws.max_column,
                }
            )
        return {"sheets": sheets}

    @staticmethod
    def _parse_ppt(file_path: str) -> dict:
        from pptx import Presentation

        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            shapes = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if text.strip():
                        shapes.append(
                            {"type": "text", "text": text, "name": shape.name}
                        )
                elif shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        table_data.append([cell.text for cell in row.cells])
                    shapes.append(
                        {
                            "type": "table",
                            "data": table_data,
                            "name": shape.name,
                        }
                    )
            slides.append({"slide": i + 1, "shapes": shapes})
        return {"total_slides": len(prs.slides), "slides": slides}
