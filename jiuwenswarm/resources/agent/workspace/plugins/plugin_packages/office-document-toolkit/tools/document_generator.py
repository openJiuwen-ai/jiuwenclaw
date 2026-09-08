from pathlib import Path
from typing import NamedTuple

from openjiuwen.core.foundation.tool import Tool, ToolCard

from text_utils import (
    CJK_PDF_TO_WORD_NOTE,
    _coerce_table,
    _table_column_count,
    coerce_generator_inputs,
    collect_structured_content_text,
    contains_cjk,
    normalize_generator_content,
    validate_generator_content,
)


class _TextBoxSpec(NamedTuple):
    text: str
    left: float
    top: float
    width: float
    height: float
    size_pt: int
    bold: bool = False


class DocumentGenerator(Tool):
    """文档生成工具：根据结构化内容生成 PDF、Word、Excel、PPT 文件。"""

    def __init__(self) -> None:
        super().__init__(
            ToolCard(
                id="document_generator",
                name="document_generator",
                description=(
                    "文档生成工具：根据结构化内容生成PDF、Word、Excel、PPT文件。"
                    "当用户需要创建文档、生成报告、输出文件时调用。"
                ),
                input_params={
                    "type": "object",
                    "properties": {
                        "format": {
                            "type": "string",
                            "enum": ["pdf", "word", "excel", "ppt", "pptx", "docx", "xlsx"],
                            "description": "输出文件格式。ppt/pptx、word/docx、excel/xlsx 等价",
                        },
                        "filename": {
                            "type": "string",
                            "description": "输出文件名（不含扩展名）",
                        },
                        "content": {
                            "type": "object",
                            "description": (
                                "结构化内容。Word/PDF: title, paragraphs[], tables[]。"
                                "PPT: slides[] 每页用 title, body, tables[]；"
                                "bullets/paragraphs/subtitle 会自动并入 body，"
                                "table/{headers,rows} 会自动并入 tables[]"
                            ),
                        },
                        "output_dir": {
                            "type": "string",
                            "description": (
                                "产物输出目录的绝对路径。传当前项目目录；"
                                "用户指定了保存位置时用用户指定的目录。"
                            ),
                        },
                    },
                    "required": ["format", "filename", "content", "output_dir"],
                },
            )
        )

    async def invoke(self, inputs, **kwargs):
        parsed = coerce_generator_inputs(inputs if isinstance(inputs, dict) else {})
        fmt = parsed.get("format", "")
        filename = parsed.get("filename", "")
        content = parsed.get("content", {})
        output_dir = parsed.get("output_dir", "")

        if not all((fmt, filename, content, output_dir)):
            return {
                "success": False,
                "error": "缺少必要参数: format, filename, content, output_dir",
            }
        if not isinstance(content, dict):
            return {"success": False, "error": "content 必须是结构化对象"}

        content = normalize_generator_content(content)
        validation_error = validate_generator_content(content, fmt)
        if validation_error:
            return {"success": False, "error": validation_error}

        base_dir = Path(output_dir).expanduser()
        base_dir.mkdir(parents=True, exist_ok=True)

        ext_map = {
            "pdf": ".pdf",
            "word": ".docx",
            "excel": ".xlsx",
            "ppt": ".pptx",
        }

        requested_format = fmt
        format_note = None
        if fmt == "pdf" and contains_cjk(collect_structured_content_text(content)):
            fmt = "word"
            format_note = CJK_PDF_TO_WORD_NOTE

        ext = ext_map.get(fmt, f".{fmt}")
        output_path = base_dir / f"{filename}{ext}"

        try:
            if fmt == "pdf":
                self._generate_pdf(str(output_path), content)
            elif fmt == "word":
                self._generate_word(str(output_path), content)
            elif fmt == "excel":
                self._generate_excel(str(output_path), content)
            elif fmt == "ppt":
                self._generate_ppt(str(output_path), content)
            else:
                return {"success": False, "error": f"不支持的格式: {fmt}"}

            if not output_path.exists() or output_path.stat().st_size == 0:
                return {
                    "success": False,
                    "error": "文件生成失败，输出文件为空或不存在",
                }

            result = {
                "success": True,
                "format": fmt,
                "path": str(output_path),
                "absolute_path": str(output_path.resolve()),
                "filename": output_path.name,
                "size_bytes": output_path.stat().st_size,
                "exists": True,
            }
            if format_note:
                result["requested_format"] = requested_format
                result["note"] = format_note
            return result
        except ImportError as e:
            return {
                "success": False,
                "error": f"依赖库缺失: {e}. 请安装对应依赖后重试。",
            }
        except Exception as e:
            return {"success": False, "error": f"生成失败: {e}"}

    async def stream(self, inputs, **kwargs):
        yield await self.invoke(inputs, **kwargs)

    @staticmethod
    def _generate_pdf(file_path: str, content: dict) -> None:
        from fpdf import FPDF

        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)

        title = content.get("title", "")
        if title:
            pdf.set_font("Helvetica", "B", 16)
            pdf.multi_cell(0, 10, title)
            pdf.ln(5)

        pdf.set_font("Helvetica", "", 11)
        for para in content.get("paragraphs", []):
            text = para if isinstance(para, str) else para.get("text", "")
            if text:
                pdf.multi_cell(0, 7, text)
                pdf.ln(3)

        for table in content.get("tables", []):
            pdf.ln(5)
            data = _coerce_table(table)
            col_count = _table_column_count(data)
            if not data or col_count < 1:
                continue
            col_width = 180 / col_count
            for row in data:
                for cell in row:
                    pdf.cell(col_width, 7, str(cell)[:50], border=1)
                pdf.ln()

        pdf.output(file_path)

    @staticmethod
    def _generate_word(file_path: str, content: dict) -> None:
        from docx import Document

        doc = Document()

        title = content.get("title", "")
        if title:
            doc.add_heading(title, level=0)

        for para in content.get("paragraphs", []):
            text = para if isinstance(para, str) else para.get("text", "")
            style = (
                para.get("style", "Normal")
                if isinstance(para, dict)
                else "Normal"
            )
            if not text:
                continue
            if style.lower().startswith("heading"):
                num_part = "".join(c for c in style if c.isdigit())
                level = int(num_part) if num_part else 1
                doc.add_heading(text, level=level)
            else:
                doc.add_paragraph(text)

        for table in content.get("tables", []):
            data = _coerce_table(table)
            cols = _table_column_count(data)
            if not data or cols < 1:
                continue
            t = doc.add_table(rows=len(data), cols=cols)
            for i, row in enumerate(data):
                for j, cell in enumerate(row):
                    if j < cols:
                        t.rows[i].cells[j].text = str(cell)

        doc.save(file_path)

    @staticmethod
    def _generate_excel(file_path: str, content: dict) -> None:
        from openpyxl import Workbook

        wb = Workbook()
        sheets = content.get("sheets", [])
        if not sheets:
            ws = wb.active
            ws.title = content.get("sheet_name", "Sheet1")
            table_data = content.get("tables", [])
            if table_data:
                data = _coerce_table(table_data[0])
            else:
                data = _coerce_table(content.get("rows") or [])
            for row in data:
                ws.append(row)
        else:
            wb.remove(wb.active)
            for sheet_data in sheets:
                ws = wb.create_sheet(
                    title=sheet_data.get("sheet_name", "Sheet")
                )
                for row in _coerce_table(sheet_data):
                    ws.append(row)

        wb.save(file_path)

    @staticmethod
    def _fill_text_frame(text_frame, text: str) -> None:
        lines = [line.strip() for line in str(text).split("\n") if line.strip()]
        if not lines:
            return
        text_frame.text = lines[0]
        for line in lines[1:]:
            paragraph = text_frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0

    @staticmethod
    def _apply_cjk_font(text_frame, size_pt: int, bold: bool = False) -> None:
        from lxml import etree
        from pptx.dml.color import RGBColor
        from pptx.oxml.ns import qn
        from pptx.util import Pt

        font_name = "Microsoft YaHei"
        color = RGBColor(0x1F, 0x4E, 0x79) if bold else RGBColor(0x33, 0x33, 0x33)
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(size_pt)
                run.font.bold = bold
                run.font.color.rgb = color
                r_pr = getattr(run, "_r").get_or_add_rPr()
                for tag in ("latin", "ea", "cs"):
                    element = r_pr.find(qn(f"a:{tag}"))
                    if element is None:
                        element = etree.SubElement(r_pr, qn(f"a:{tag}"))
                    element.set("typeface", font_name)

    @staticmethod
    def _add_textbox(slide, spec: _TextBoxSpec):
        from pptx.util import Inches

        box = slide.shapes.add_textbox(
            Inches(spec.left), Inches(spec.top), Inches(spec.width), Inches(spec.height)
        )
        text_frame = box.text_frame
        text_frame.word_wrap = True
        DocumentGenerator._fill_text_frame(text_frame, spec.text)
        DocumentGenerator._apply_cjk_font(text_frame, size_pt=spec.size_pt, bold=spec.bold)
        return box

    @staticmethod
    def _add_slide_table(slide, data, top_inches: float = 1.6) -> float:
        from pptx.util import Inches

        rows = len(data)
        cols = _table_column_count(data)
        if rows < 1 or cols < 1:
            return top_inches
        height = min(4.8, max(0.6, 0.32 * rows + 0.2))
        table = slide.shapes.add_table(
            rows, cols, Inches(0.5), Inches(top_inches), Inches(9.0), Inches(height)
        ).table
        for i, row in enumerate(data):
            for j, cell in enumerate(row):
                if j < cols:
                    table.cell(i, j).text = str(cell)
                    DocumentGenerator._apply_cjk_font(
                        table.cell(i, j).text_frame,
                        size_pt=12,
                        bold=(i == 0),
                    )
        return top_inches + height + 0.2

    @staticmethod
    def _generate_ppt(file_path: str, content: dict) -> None:
        from pptx import Presentation

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slides = content.get("slides") or []
        if not slides and content.get("title"):
            slides = [
                {
                    "title": content.get("title", ""),
                    "body": content.get("subtitle", ""),
                }
            ]

        for slide_data in slides:
            if not isinstance(slide_data, dict):
                continue
            body = slide_data.get("body") or ""
            if not isinstance(body, str):
                body = "\n".join(str(item) for item in body if item)
            tables = [_coerce_table(item) for item in slide_data.get("tables", [])]
            tables = [item for item in tables if _table_column_count(item) > 0]
            slide = prs.slides.add_slide(blank_layout)
            slide_title = slide_data.get("title", "")
            if slide_title:
                DocumentGenerator._add_textbox(
                    slide,
                    _TextBoxSpec(
                        str(slide_title), 0.5, 0.25, 9.0, 0.8, 28, bold=True
                    ),
                )
            if body:
                DocumentGenerator._add_textbox(
                    slide,
                    _TextBoxSpec(body, 0.5, 1.15, 9.0, 2.0 if tables else 5.8, 18),
                )
            table_top = 3.3 if body else 1.2
            for data in tables:
                table_top = DocumentGenerator._add_slide_table(
                    slide, data, top_inches=table_top
                )

        if not prs.slides:
            raise ValueError("PPT 没有有效幻灯片，请提供 slides[].title 与 body/tables")
        prs.save(file_path)
