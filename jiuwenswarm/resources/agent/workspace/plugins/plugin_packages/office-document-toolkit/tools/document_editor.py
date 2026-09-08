import os

from openjiuwen.core.foundation.tool import Tool, ToolCard


class DocumentEditor(Tool):
    """文档编辑工具：编辑已有 Word、Excel、PPT 文档（非从零新建）。"""

    def __init__(self) -> None:
        super().__init__(
            ToolCard(
                id="document_editor",
                name="document_editor",
                description=(
                    "文档编辑工具：在现有Word里替换文本/追加段落/插入图片、"
                    "往现有Excel指定单元格写数、用母版模板套出PPT。"
                    "当用户需要修改已有文档而非新建时调用。"
                ),
                input_params={
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "要编辑的文件绝对路径",
                        },
                        "format": {
                            "type": "string",
                            "enum": ["word", "excel", "ppt"],
                            "description": "文档格式",
                        },
                        "operations": {
                            "type": "array",
                            "items": {"type": "object"},
                            "description": (
                                "操作列表，每项含 type 和参数。"
                                "Word: replace_text/append_paragraph/"
                                "insert_heading/insert_image/add_table/"
                                "accept_all_changes/reject_all_changes/"
                                "track_changes。"
                                "Excel: write_cell/write_range/append_row/"
                                "add_formula。"
                                "PPT: replace_placeholder/"
                                "add_slide_from_template/replace_slide_text"
                            ),
                        },
                    },
                    "required": ["file_path", "format", "operations"],
                },
            )
        )

    async def invoke(self, inputs, **kwargs):
        file_path = inputs.get("file_path", "")
        fmt = inputs.get("format", "")
        operations = inputs.get("operations", [])

        if not file_path or not os.path.isfile(file_path):
            return {"success": False, "error": f"文件不存在: {file_path}"}
        if not fmt:
            return {"success": False, "error": "缺少 format 参数"}
        if not operations:
            return {"success": False, "error": "缺少 operations 参数"}

        try:
            if fmt == "word":
                results = self._edit_word(file_path, operations)
            elif fmt == "excel":
                results = self._edit_excel(file_path, operations)
            elif fmt == "ppt":
                results = self._edit_ppt(file_path, operations)
            else:
                return {"success": False, "error": f"不支持的格式: {fmt}"}

            return {
                "success": True,
                "format": fmt,
                "file_path": file_path,
                "operations_executed": len(results),
                "results": results,
                "size_bytes": os.path.getsize(file_path),
                "exists": True,
            }
        except ImportError as e:
            return {
                "success": False,
                "error": f"依赖库缺失: {e}. 请安装对应依赖后重试。",
            }
        except Exception as e:
            return {"success": False, "error": f"编辑失败: {e}"}

    async def stream(self, inputs, **kwargs):
        yield await self.invoke(inputs, **kwargs)

    @staticmethod
    def _edit_word(file_path: str, operations: list) -> list:
        from docx import Document
        from docx.shared import Inches

        doc = Document(file_path)
        results = []

        for op in operations:
            op_type = op.get("type", "")

            if op_type == "replace_text":
                old_text = op.get("old_text", "")
                new_text = op.get("new_text", "")
                count = 0
                for para in doc.paragraphs:
                    if old_text in para.text:
                        for run in para.runs:
                            if old_text in run.text:
                                run.text = run.text.replace(
                                    old_text, new_text
                                )
                                count += 1
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for para in cell.paragraphs:
                                if old_text in para.text:
                                    for run in para.runs:
                                        if old_text in run.text:
                                            run.text = run.text.replace(
                                                old_text, new_text
                                            )
                                            count += 1
                results.append(
                    {"type": "replace_text", "replacements": count}
                )

            elif op_type == "append_paragraph":
                text = op.get("text", "")
                style = op.get("style", "Normal")
                doc.add_paragraph(text, style=style)
                results.append(
                    {"type": "append_paragraph", "text": text[:50]}
                )

            elif op_type == "insert_heading":
                text = op.get("text", "")
                level = op.get("level", 1)
                doc.add_heading(text, level=level)
                results.append(
                    {
                        "type": "insert_heading",
                        "text": text[:50],
                        "level": level,
                    }
                )

            elif op_type == "insert_image":
                image_path = op.get("image_path", "")
                width = op.get("width")
                if image_path and os.path.isfile(image_path):
                    pic_kwargs = {}
                    if width:
                        pic_kwargs["width"] = Inches(width)
                    doc.add_picture(image_path, **pic_kwargs)
                    results.append(
                        {"type": "insert_image", "image": image_path}
                    )
                else:
                    results.append(
                        {
                            "type": "insert_image",
                            "error": "图片文件不存在",
                        }
                    )

            elif op_type == "add_table":
                data = op.get("data", [])
                if data:
                    rows = len(data)
                    cols = max(len(r) for r in data) if data else 1
                    table = doc.add_table(rows=rows, cols=cols)
                    for i, row in enumerate(data):
                        for j, cell in enumerate(row):
                            if j < cols:
                                table.rows[i].cells[j].text = str(cell)
                    results.append(
                        {"type": "add_table", "rows": rows, "cols": cols}
                    )

            elif op_type == "accept_all_changes":
                nsmap = {
                    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                }
                body = doc.element.body
                for ins in body.findall(".//w:ins", nsmap):
                    parent = ins.getparent()
                    for child in list(ins):
                        parent.insert(list(parent).index(ins), child)
                    parent.remove(ins)
                for dele in body.findall(".//w:del", nsmap):
                    dele.getparent().remove(dele)
                results.append(
                    {"type": "accept_all_changes", "status": "已接受所有修订"}
                )

            elif op_type == "reject_all_changes":
                nsmap = {
                    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                }
                body = doc.element.body
                for ins in body.findall(".//w:ins", nsmap):
                    ins.getparent().remove(ins)
                for dele in body.findall(".//w:del", nsmap):
                    parent = dele.getparent()
                    for child in list(dele):
                        parent.insert(list(parent).index(dele), child)
                    parent.remove(dele)
                results.append(
                    {"type": "reject_all_changes", "status": "已拒绝所有修订"}
                )

            elif op_type == "track_changes":
                from docx.oxml import OxmlElement
                from docx.oxml.ns import qn

                enable = op.get("enable", True)
                settings = doc.settings.element
                existing = settings.find(qn("w:trackChanges"))
                if enable:
                    if existing is None:
                        tc = OxmlElement("w:trackChanges")
                        settings.append(tc)
                else:
                    if existing is not None:
                        settings.remove(existing)
                results.append(
                    {"type": "track_changes", "enabled": enable}
                )

        doc.save(file_path)
        return results

    @staticmethod
    def _edit_excel(file_path: str, operations: list) -> list:
        from openpyxl import load_workbook
        from openpyxl.utils import range_to_tuple

        wb = load_workbook(file_path)
        results = []

        for op in operations:
            op_type = op.get("type", "")

            if op_type == "write_cell":
                sheet_name = op.get("sheet_name")
                cell_ref = op.get("cell_ref", "")
                value = op.get("value", "")
                ws = wb[sheet_name] if sheet_name else wb.active
                ws[cell_ref] = value
                results.append(
                    {
                        "type": "write_cell",
                        "cell": cell_ref,
                        "value": str(value)[:50],
                    }
                )

            elif op_type == "write_range":
                sheet_name = op.get("sheet_name")
                start_cell = op.get("start_cell", "A1")
                data = op.get("data", [])
                ws = wb[sheet_name] if sheet_name else wb.active
                start_row, start_col = range_to_tuple(start_cell)
                for i, row in enumerate(data):
                    for j, val in enumerate(row):
                        ws.cell(
                            row=start_row + i,
                            column=start_col + j,
                            value=val,
                        )
                results.append(
                    {
                        "type": "write_range",
                        "start_cell": start_cell,
                        "rows": len(data),
                    }
                )

            elif op_type == "append_row":
                sheet_name = op.get("sheet_name")
                data = op.get("data", [])
                ws = wb[sheet_name] if sheet_name else wb.active
                ws.append(data)
                results.append(
                    {"type": "append_row", "columns": len(data)}
                )

            elif op_type == "add_formula":
                sheet_name = op.get("sheet_name")
                cell_ref = op.get("cell_ref", "")
                formula = op.get("formula", "")
                ws = wb[sheet_name] if sheet_name else wb.active
                ws[cell_ref] = formula
                results.append(
                    {
                        "type": "add_formula",
                        "cell": cell_ref,
                        "formula": formula,
                    }
                )

        wb.save(file_path)
        return results

    @staticmethod
    def _edit_ppt(file_path: str, operations: list) -> list:
        from pptx import Presentation

        prs = Presentation(file_path)
        results = []

        for op in operations:
            op_type = op.get("type", "")

            if op_type == "replace_placeholder":
                placeholder = op.get("placeholder", "")
                new_text = op.get("new_text", "")
                count = 0
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            if placeholder in shape.text_frame.text:
                                for para in (
                                    shape.text_frame.paragraphs
                                ):
                                    for run in para.runs:
                                        if placeholder in run.text:
                                            run.text = run.text.replace(
                                                placeholder, new_text
                                            )
                                            count += 1
                results.append(
                    {
                        "type": "replace_placeholder",
                        "replacements": count,
                    }
                )

            elif op_type == "add_slide_from_template":
                layout_index = op.get("layout_index", 1)
                title = op.get("title", "")
                body = op.get("body", "")
                if layout_index < len(prs.slide_layouts):
                    slide_layout = prs.slide_layouts[layout_index]
                    slide = prs.slides.add_slide(slide_layout)
                    if title and slide.shapes.title:
                        slide.shapes.title.text = title
                    if body and len(slide.placeholders) > 1:
                        slide.placeholders[1].text_frame.text = body
                    results.append(
                        {
                            "type": "add_slide_from_template",
                            "title": title[:50],
                        }
                    )
                else:
                    results.append(
                        {
                            "type": "add_slide_from_template",
                            "error": "布局索引超出范围",
                        }
                    )

            elif op_type == "replace_slide_text":
                slide_index = op.get("slide_index", 0)
                old_text = op.get("old_text", "")
                new_text = op.get("new_text", "")
                if slide_index < len(prs.slides):
                    slide = prs.slides[slide_index]
                    count = 0
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in (
                                shape.text_frame.paragraphs
                            ):
                                for run in para.runs:
                                    if old_text in run.text:
                                        run.text = run.text.replace(
                                            old_text, new_text
                                        )
                                        count += 1
                    results.append(
                        {
                            "type": "replace_slide_text",
                            "slide": slide_index,
                            "replacements": count,
                        }
                    )

        prs.save(file_path)
        return results
